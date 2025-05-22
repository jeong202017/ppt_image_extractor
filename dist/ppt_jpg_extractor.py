import os
import shutil
import re
import threading
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tkinter import Tk, Label, Button, filedialog, messagebox

def sanitize_filename(name):
    """파일 이름으로 안전한 문자열 생성"""
    name = re.sub(r'[\\/:"*?<>|]', '_', name)
    return name.replace('.', '_')

def remove_images_and_save(prs, output_path):
    """PPT에서 이미지 제거하고 저장"""
    for slide in prs.slides:
        for shape in list(slide.shapes):  # 복사본으로 순회
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide.shapes._spTree.remove(shape._element)
    prs.save(output_path)

def extract_images_from_all_pptx(input_dir, output_dir_images, output_dir_no_images):
    os.makedirs(output_dir_images, exist_ok=True)
    os.makedirs(output_dir_no_images, exist_ok=True)

    pptx_files = [f for f in os.listdir(input_dir) if f.endswith('.pptx')]
    if not pptx_files:
        messagebox.showwarning("경고", "선택한 폴더에 pptx 파일이 없습니다.")
        return

    for filename in pptx_files:
        pptx_path = os.path.join(input_dir, filename)
        prs = Presentation(pptx_path)
        image_count = 0

        ppt_name = sanitize_filename(os.path.splitext(filename)[0])
        ppt_image_folder = os.path.join(output_dir_images, ppt_name)
        os.makedirs(ppt_image_folder, exist_ok=True)

        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    content_type = image.content_type
                    ext = content_type.split('/')[-1]
                    image_filename = f"slide{slide_index+1}_img{shape_index+1}.{ext}"
                    image_path = os.path.join(ppt_image_folder, image_filename)

                    try:
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        image_count += 1
                    except OSError as e:
                        print(f"[⚠️] 이미지 저장 실패: {image_path}\n{e}")

        if image_count > 0:
            # 이미지 제거한 PPT 저장
            output_clean_pptx = os.path.join(output_dir_no_images, filename)
            remove_images_and_save(prs, output_clean_pptx)
            print(f"[✅] {filename} → 이미지 {image_count}개 추출 + 텍스트 전용 PPT 저장")
        else:
            shutil.copy(pptx_path, os.path.join(output_dir_no_images, filename))
            shutil.rmtree(ppt_image_folder)
            print(f"[❌] {filename} → 이미지 없음 → 그대로 복사")

    messagebox.showinfo("완료", "이미지 추출 및 텍스트 전용 PPT 저장 완료")

def start_extraction():
    input_dir = filedialog.askdirectory(title="PPTX 폴더 선택")
    if not input_dir:
        return

    output_dir_images = os.path.join(input_dir, 'extracted_images')
    output_dir_no_images = os.path.join(input_dir, 'no_image_pptx')

    threading.Thread(
        target=extract_images_from_all_pptx,
        args=(input_dir, output_dir_images, output_dir_no_images)
    ).start()

def run_gui():
    root = Tk()
    root.title("PPTX 이미지 추출 및 제거기")
    root.geometry("420x200")

    Label(root, text="📁 폴더 내 PPTX 파일에서 이미지를 추출하고,\n텍스트만 남긴 PPT를 저장합니다.", pady=20).pack()
    Button(root, text="PPTX 폴더 선택 및 실행", command=start_extraction, padx=20, pady=10).pack()
    Label(root, text="© 2025 PPT Extractor", pady=20).pack()

    root.mainloop()

if __name__ == "__main__":
    run_gui()
